"""
精益文档知识库构建脚本（离线版）
使用 TF-IDF 向量化，不需要下载任何模型，全程离线运行
"""

import os
import sys
import json
import pickle
import hashlib
import tempfile
import subprocess
from pathlib import Path
from typing import List, Dict, Tuple

# 文档解析库
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader

# TF-IDF 向量化
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np

# ====== 配置 ======
# 修改为你的精益文档所在目录
# 例如: DOCUMENTS_DIR = Path(r"C:\Users\YourName\Desktop\精益文档")
DOCUMENTS_DIR = Path(r"C:\Users\YourName\Desktop\精益文档")
OUTPUT_DIR = Path(__file__).parent / "knowledge_base"

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".doc", ".txt", ".xlsx", ".md"}


# ====== 文档解析器 ======

def parse_pdf(file_path: Path) -> str:
    text_parts = []
    try:
        reader = PdfReader(str(file_path))
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
    except Exception as e:
        print(f"  [!] PDF解析失败 {file_path.name}: {e}")
    return "\n".join(text_parts)


def parse_pptx(file_path: Path) -> str:
    text_parts = []
    try:
        prs = Presentation(str(file_path))
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(row_texts))
            if slide_texts:
                text_parts.append(f"【第{i}页】\n" + "\n".join(slide_texts))
    except Exception as e:
        print(f"  [!] PPTX解析失败 {file_path.name}: {e}")
    return "\n".join(text_parts)


def parse_docx(file_path: Path) -> str:
    text_parts = []
    try:
        doc = Document(str(file_path))
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells]
                text_parts.append(" | ".join(row_texts))
    except Exception as e:
        print(f"  [!] DOCX解析失败 {file_path.name}: {e}")
    return "\n".join(text_parts)


def parse_txt(file_path: Path) -> str:
    encodings = ["utf-8", "gbk", "gb2312", "utf-16"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    print(f"  [!] 无法解码 {file_path.name}")
    return ""


def parse_xlsx(file_path: Path) -> str:
    text_parts = []
    try:
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
        for sheet in wb.worksheets:
            sheet_name = sheet.title
            rows_text = []
            for row in sheet.iter_rows():
                row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                row_line = " | ".join(row_values).strip()
                if row_line:
                    rows_text.append(row_line)
            if rows_text:
                text_parts.append(f"【工作表: {sheet_name}】")
                text_parts.extend(rows_text)
        wb.close()
    except Exception as e:
        print(f"  [!] XLSX解析失败 {file_path.name}: {e}")
    return "\n".join(text_parts)


def parse_doc(file_path: Path) -> str:
    """解析旧版 .doc 文件（使用 Word COM 自动化）"""
    try:
        # 尝试使用 Word COM
        import win32com.client
        tmp_txt = tempfile.mktemp(suffix=".txt")
        try:
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            # 打开 .doc 文件
            doc = word.Documents.Open(str(file_path))
            # 另存为纯文本 (wdFormatText=2)
            doc.SaveAs(tmp_txt, FileFormat=2)
            doc.Close()
            word.Quit()

            # 读取保存的文本（通常为 gbk 编码）
            for enc in ["gbk", "gb2312", "utf-8", "utf-16-le"]:
                try:
                    with open(tmp_txt, "r", encoding=enc) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            return ""
        except Exception as e:
            print(f"  [!] Word COM 失败: {e}")
            return ""
        finally:
            if os.path.exists(tmp_txt):
                try:
                    os.remove(tmp_txt)
                except:
                    pass
    except ImportError:
        print(f"  [!] pywin32 未安装，跳过 .doc 文件")
        return ""


PARSERS = {
    ".pdf": parse_pdf,
    ".pptx": parse_pptx,
    ".docx": parse_docx,
    ".doc": parse_doc,
    ".txt": parse_txt,
    ".md": parse_txt,
    ".xlsx": parse_xlsx,
}


def chunk_text(text: str, file_name: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict]:
    """将长文本切分为多个块"""
    if not text.strip():
        return []

    # 按句子切分
    separators = ["\n\n", "\n", "。", "！", "？", "；", "。\n", "！\n", "？\n"]
    sentences = [text]
    for sep in separators:
        new_sentences = []
        for s in sentences:
            if len(s) > chunk_size:
                new_sentences.extend(s.split(sep))
            else:
                new_sentences.append(s)
        sentences = [s.strip() for s in new_sentences if s.strip()]

    # 合并成块
    chunks = []
    current_chunk = ""
    chunk_id = 0

    for sentence in sentences:
        if len(current_chunk) + len(sentence) <= chunk_size:
            current_chunk += sentence
        else:
            if current_chunk:
                chunks.append({
                    "id": f"{file_name}_{chunk_id}",
                    "text": current_chunk,
                    "source": file_name,
                })
                chunk_id += 1
                # 保留 overlap
                if overlap > 0:
                    words = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                    current_chunk = words + sentence
                else:
                    current_chunk = sentence

    if current_chunk:
        chunks.append({
            "id": f"{file_name}_{chunk_id}",
            "text": current_chunk,
            "source": file_name,
        })

    return chunks


def load_and_chunk_documents(directory: Path) -> Tuple[List[Dict], List[str]]:
    """加载所有文档并切分为块"""
    all_chunks = []
    file_count = 0
    success_count = 0

    for file_path in sorted(directory.rglob("*")):
        if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        if file_path.name.startswith("~$"):
            continue

        file_count += 1
        rel_path = str(file_path.relative_to(directory))
        print(f"  处理: {rel_path}")

        parser = PARSERS.get(file_path.suffix.lower())
        if not parser:
            continue

        text = parser(file_path)
        if text.strip():
            chunks = chunk_text(text, rel_path)
            all_chunks.extend(chunks)
            success_count += 1
            print(f"    -> 提取 {len(chunks)} 个文本块 ({len(text)} 字符)")
        else:
            print(f"    [!] 未提取到文本内容")

    print(f"\n扫描完成: {file_count} 个文件, 成功 {success_count} 个")
    print(f"共 {len(all_chunks)} 个文本块")
    return all_chunks


def build_knowledge_base():
    """主流程"""
    import time
    start = time.time()

    print("=" * 60)
    print("  精益文档知识库构建工具 (离线版)")
    print("=" * 60)

    # Step 1: 创建输出目录
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Step 2: 加载并切分文档
    print(f"\n[扫描] {DOCUMENTS_DIR}")
    chunks = load_and_chunk_documents(DOCUMENTS_DIR)

    if not chunks:
        print("\n[失败] 未提取到任何文档内容")
        return

    # Step 3: 构建 TF-IDF 向量索引
    print("\n[索引] 构建 TF-IDF 向量索引...")
    texts = [c["text"] for c in chunks]

    vectorizer = TfidfVectorizer(
        max_features=50000,
        ngram_range=(1, 2),
        analyzer="char_wb",  # 对中文友好
        sublinear_tf=True,
    )

    tfidf_matrix = vectorizer.fit_transform(texts)
    print(f"   词汇表大小: {len(vectorizer.get_feature_names_out())}")
    print(f"   矩阵形状: {tfidf_matrix.shape}")

    # Step 4: 保存到磁盘
    print(f"\n[保存] 保存到 {OUTPUT_DIR}")

    with open(OUTPUT_DIR / "chunks.pkl", "wb") as f:
        pickle.dump(chunks, f)

    with open(OUTPUT_DIR / "vectorizer.pkl", "wb") as f:
        pickle.dump(vectorizer, f)

    with open(OUTPUT_DIR / "tfidf_matrix.pkl", "wb") as f:
        pickle.dump(tfidf_matrix, f)

    # 保存文件列表供参考
    sources = sorted(set(c["source"] for c in chunks))
    with open(OUTPUT_DIR / "sources.txt", "w", encoding="utf-8") as f:
        for s in sources:
            f.write(s + "\n")

    elapsed = time.time() - start
    print("\n" + "=" * 60)
    print("  [完成] 知识库构建完成!")
    print(f"  [统计] 共 {len(chunks)} 个文本块")
    print(f"  [来源] {len(sources)} 个文件")
    print(f"  [耗时] {elapsed:.1f} 秒")
    print(f"  [路径] {OUTPUT_DIR}")
    print("=" * 60)


if __name__ == "__main__":
    build_knowledge_base()
